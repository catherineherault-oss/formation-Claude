#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère le diaporama de copil (PPTX) à partir des résultats de l'étude vente directe.
Sortie : outputs/presentations/copil_vente_directe.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures/png"
OUT = ROOT / "outputs/presentations/copil_vente_directe.pptx"

# ---- Palette --------------------------------------------------------------
VERT   = RGBColor(0x2E, 0x5E, 0x3A)   # vert foncé (titre)
VERTC  = RGBColor(0x4C, 0x8C, 0x5A)   # vert clair (accent)
ROUGE  = RGBColor(0xC0, 0x39, 0x2B)   # rouge (chiffre-choc)
GRIS   = RGBColor(0x44, 0x44, 0x44)
BLANC  = RGBColor(0xFF, 0xFF, 0xFF)
CLAIR  = RGBColor(0xF2, 0xF6, 0xF3)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(s, x, y, w, h, text, size=18, color=GRIS, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Calibri", line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb

def bullets(s, x, y, w, h, items, size=20, color=GRIS, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, (txt_, b) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        r = p.add_run(); r.text = "•  " + txt_
        r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb

def header(s, num, title):
    rect(s, 0, 0, W, Inches(1.15), VERT)
    rect(s, 0, Inches(1.15), W, Inches(0.06), VERTC)
    txt(s, Inches(0.5), Inches(0.18), W - Inches(1), Inches(0.85),
        (f"{num}.  " if num else "") + title, size=30, color=BLANC, bold=True,
        anchor=MSO_ANCHOR.MIDDLE)

def footer(s, page):
    txt(s, Inches(0.5), H - Inches(0.45), Inches(9), Inches(0.35),
        "Vente directe alimentaire — Comité de pilotage — Enquête 1 025 ménages",
        size=10, color=RGBColor(0x99,0x99,0x99))
    txt(s, W - Inches(1.2), H - Inches(0.45), Inches(0.8), Inches(0.35),
        str(page), size=10, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)

def chip(s, x, y, w, number, label, color=ROUGE, hh=1.7):
    rect(s, x, y, w, Inches(hh), CLAIR)
    txt(s, x, y + Inches(0.2), w, Inches(0.85), number, size=44, color=color,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.15), y + Inches(1.05), w - Inches(0.3), Inches(0.6), label,
        size=13, color=GRIS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ====================================================================== 1 TITRE
s = slide()
rect(s, 0, 0, W, H, VERT)
rect(s, 0, Inches(3.05), W, Inches(0.08), VERTC)
txt(s, Inches(1), Inches(2.0), W - Inches(2), Inches(1.1),
    "L'achat alimentaire en vente directe", size=44, color=BLANC, bold=True)
txt(s, Inches(1), Inches(3.25), W - Inches(2), Inches(1.2),
    "Quelle place dans les stratégies d'approvisionnement\ndes ménages français ?",
    size=24, color=RGBColor(0xD8,0xE8,0xDC))
txt(s, Inches(1), Inches(5.6), W - Inches(2), Inches(0.9),
    "Enquête représentative nationale · 1 025 ménages · 13 canaux d'achat\nComité de pilotage — juillet 2026",
    size=15, color=RGBColor(0xB9,0xD4,0xC0))

# ============================================================ 2 MESSAGE CENTRAL
s = slide()
rect(s, 0, 0, W, H, CLAIR)
rect(s, 0, Inches(2.4), Inches(0.25), Inches(2.7), ROUGE)
txt(s, Inches(0.8), Inches(2.3), W - Inches(1.6), Inches(1.6),
    "La vente directe est partout…\nmais pèse presque rien.",
    size=40, color=VERT, bold=True)
txt(s, Inches(0.85), Inches(4.5), W - Inches(1.7), Inches(1.2),
    "2 ménages sur 3 y ont recours — mais elle reste un complément d'appoint, "
    "pas une alternative à la grande distribution.",
    size=22, color=GRIS)
footer(s, 2)

# ===================================================== 3 PARADOXE (fig1)
s = slide(); header(s, 1, "Un paradoxe d'ampleur")
chip(s, Inches(0.5), Inches(1.5), Inches(3.5), "65 %", "des ménages achètent\nen vente directe", ROUGE)
chip(s, Inches(0.5), Inches(3.3), Inches(3.5), "1,5 %", "seulement du budget\nalimentaire", ROUGE)
chip(s, Inches(0.5), Inches(5.1), Inches(3.5), "58 %", "captés par le seul\nhypermarché", VERT)
if (FIG/"fig1_part_budget_canal.png").exists():
    s.shapes.add_picture(str(FIG/"fig1_part_budget_canal.png"), Inches(4.4), Inches(1.45),
                         height=Inches(5.1))
txt(s, Inches(4.4), Inches(6.75), Inches(8.4), Inches(0.5),
    "→ Un canal de complémentarité ciblée, non de substitution.",
    size=15, color=VERT, bold=True)
footer(s, 3)

# ===================================================== 4 ESSAI vs ADOPTION
s = slide(); header(s, 2, "Essayer n'est pas adopter")
chip(s, Inches(1.2), Inches(2.4), Inches(4.2), "65 %", "ont DÉJÀ acheté\nen vente directe", VERTC)
txt(s, Inches(5.6), Inches(2.7), Inches(1.6), Inches(1.5), "➜", size=54,
    color=GRIS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
chip(s, Inches(7.4), Inches(2.4), Inches(4.2), "17 %", "achètent RÉGULIÈREMENT\n(au moins 1×/mois)", ROUGE)
txt(s, Inches(1.2), Inches(4.9), Inches(10.9), Inches(1.6),
    "Beaucoup essaient, peu installent la vente directe dans leurs habitudes.\n"
    "Le vrai plafond de verre n'est pas l'adhésion, c'est LA RÉGULARITÉ.",
    size=22, color=VERT, bold=True)
footer(s, 4)

# ===================================================== 5 TRANSVERSALITE
s = slide(); header(s, 3, "Ce n'est pas une pratique de « niche »")
bullets(s, Inches(0.7), Inches(1.7), Inches(12), Inches(4), [
    ("Le profil social explique TRÈS PEU qui achète en direct (modèle quasi sans pouvoir prédictif).", True),
    ("Les cadres y recourent un peu plus — mais l'écart reste modéré.", False),
    ("La pratique TRAVERSE toutes les catégories sociales et tous les âges.", False),
    ("Il existe même un groupe de ménages à PETIT BUDGET qui se fournit hors hypermarché, sur les marchés.", False),
], size=22, gap=14)
rect(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.1), CLAIR)
txt(s, Inches(0.95), Inches(5.85), Inches(11.4), Inches(0.9),
    "→ Implication : cibler large. La vente directe n'est pas socialement clivante ; "
    "les dispositifs peuvent viser tous les publics, budgets modestes compris.",
    size=17, color=VERT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# ===================================================== 6 TYPOLOGIE (table)
s = slide(); header(s, 4, "Cinq stratégies d'approvisionnement")
rows = [
    ("Stratégie", "Part", "Ont essayé\nla VD", "VD\nrégulière"),
    ("Généralistes (hyper + compléments)", "61 %", "72 %", "16 %"),
    ("Captifs de la grande distribution", "24 %", "28 %", "6 %"),
    ("Omnivores engagés (vrais réguliers)", "8 %", "97 %", "67 %"),
    ("Multi-canal conventionnels (essaient tout)", "5 %", "100 %", "19 %"),
    ("Hors-hypermarché (petits budgets)", "2,5 %", "81 %", "8 %"),
]
tw = Inches(11.8); th = Inches(4.4)
tbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(1.6), tw, th).table
tbl.columns[0].width = Inches(6.4)
for c in (1,2,3): tbl.columns[c].width = Inches(1.8)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c); cell.text = ""
        p = cell.text_frame.paragraphs[0]; run = p.add_run(); run.text = val
        run.font.size = Pt(15 if r else 13); run.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        if r == 0:
            run.font.bold = True; run.font.color.rgb = BLANC
            cell.fill.solid(); cell.fill.fore_color.rgb = VERT
        else:
            run.font.color.rgb = GRIS
            cell.fill.solid(); cell.fill.fore_color.rgb = BLANC if r % 2 else CLAIR
            if c == 3:  # colonne VD régulière en gras rouge
                run.font.bold = True; run.font.color.rgb = ROUGE
txt(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.7),
    "Un seul groupe (les « omnivores engagés ») fait vraiment de la vente directe une habitude.",
    size=15, color=VERT, bold=True)
footer(s, 6)

# ===================================================== 7 MOTIVATIONS (fig5)
s = slide(); header(s, 5, "On y va pour le local, pas pour le prix")
if (FIG/"fig5_motivations_marche.png").exists():
    s.shapes.add_picture(str(FIG/"fig5_motivations_marche.png"), Inches(0.5), Inches(1.5),
                         height=Inches(5.0))
bullets(s, Inches(7.6), Inches(1.9), Inches(5.4), Inches(4), [
    ("Produits locaux — 51 %", True),
    ("Circuit court — 41 %", False),
    ("Goût des produits — 39 %", False),
    ("Soutien aux agriculteurs — 34 %", False),
    ("Le PRIX seulement 5ᵉ — 18 %", True),
], size=20, gap=12)
txt(s, Inches(7.6), Inches(5.9), Inches(5.4), Inches(1.1),
    "→ Communiquer « local & qualité »,\nnon « pas cher ».",
    size=18, color=ROUGE, bold=True)
footer(s, 7)

# ===================================================== 8 LECTURE STRATEGIQUE
s = slide(); header(s, None, "Ce que ça veut dire")
cards = [
    ("Le potentiel de conquête\nest déjà atteint",
     "65 % ont déjà essayé. Inutile de « convertir » — le levier est de faire passer de l'essai à l'habitude."),
    ("Les freins sont pratiques,\npas idéologiques",
     "Temps, accessibilité, régularité et diversité de l'offre. Pas un problème d'adhésion."),
    ("Cibler large,\npas une niche",
     "La pratique n'est pas socialement clivante : viser tous les publics, budgets modestes compris."),
]
x = Inches(0.6)
for i,(t,d) in enumerate(cards):
    cx = x + i*Inches(4.15)
    rect(s, cx, Inches(1.7), Inches(3.9), Inches(4.6), CLAIR)
    rect(s, cx, Inches(1.7), Inches(3.9), Inches(0.12), VERTC)
    txt(s, cx+Inches(0.25), Inches(2.0), Inches(3.4), Inches(1.3), t, size=20,
        color=VERT, bold=True)
    txt(s, cx+Inches(0.25), Inches(3.5), Inches(3.4), Inches(2.6), d, size=16,
        color=GRIS, line_spacing=1.1)
footer(s, 8)

# ===================================================== 9 RECOMMANDATIONS
s = slide(); header(s, None, "3 leviers d'action")
recs = [
    ("1", "Miser sur la fidélisation, pas l'acquisition",
     "Transformer l'essai en habitude : abonnements, points de retrait réguliers, élargissement des gammes."),
    ("2", "Lever les freins logistiques",
     "Horaires, proximité, régularité de l'offre — c'est là que se joue le passage à l'usage régulier."),
    ("3", "Communiquer « local & qualité », pas « pas cher »",
     "C'est le registre qui mobilise réellement les ménages."),
]
y = Inches(1.6)
for num, t, d in recs:
    rect(s, Inches(0.7), y, Inches(1.0), Inches(1.55), VERT)
    txt(s, Inches(0.7), y, Inches(1.0), Inches(1.55), num, size=40, color=BLANC,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.0), y+Inches(0.12), Inches(10.6), Inches(0.6), t, size=21,
        color=VERT, bold=True)
    txt(s, Inches(2.0), y+Inches(0.78), Inches(10.6), Inches(0.8), d, size=16, color=GRIS)
    y += Inches(1.8)
footer(s, 9)

# ===================================================== 10 METHODE (annexe)
s = slide(); header(s, None, "Méthode & portée")
bullets(s, Inches(0.7), Inches(1.7), Inches(12), Inches(3.6), [
    ("Enquête par questionnaire auprès d'un échantillon représentatif national de 1 025 ménages (quotas sexe, âge, CSP, région).", False),
    ("13 canaux d'achat mesurés : fréquence et part de budget de chacun.", False),
    ("Typologie par analyse des correspondances multiples + classification ; régression sur les déterminants.", False),
    ("« Vente directe » = achat direct au producteur (à la ferme, AMAP, magasins de producteurs, marchés).", False),
], size=18, gap=12)
rect(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.5), CLAIR)
txt(s, Inches(0.95), Inches(5.55), Inches(11.4), Inches(1.2),
    "Limites : enquête transversale (un instantané, pas une évolution dans le temps) ; "
    "données déclaratives (fréquences et budgets estimés par les répondants).",
    size=15, color=GRIS, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10)

prs.save(str(OUT))
print("PPTX généré :", OUT, "|", len(prs.slides._sldIdLst), "diapositives")
